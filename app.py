from flask import Flask, request, send_file, render_template
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)

# Basic HTML for the upload form (can be served from render_template)
@app.route('/')
def index():
    return render_template('index.html') # Assuming the HTML above is in templates/index.html

@app.route('/upload-and-design', methods=['POST'])
def upload_and_design():
    if 'pptFile' not in request.files:
        return "No file part", 400

    file = request.files['pptFile']
    if file.filename == '':
        return "No selected file", 400

    if file:
        try:
            # Read the incoming PPTX file into a BytesIO object
            prs = Presentation(file)

            # --- Design Logic ---
            # 1. Accessing existing slides
            # You might iterate through slides to find "blank" ones.
            # A blank slide might have only a title and content placeholder,
            # or even fewer elements. This detection can be complex.
            # For simplicity, let's just apply design to ALL slides.

            for slide_idx, slide in enumerate(prs.slides):
                # Example 1: Add a simple text box
                left = top = width = height = prs.slide_width / 4 # Example dimensions
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                p = tf.add_paragraph()
                p.text = f"Design Trace {slide_idx + 1}"
                p.font.size = 24  # Set font size
                p.font.bold = True
                p.text_color.rgb = (0xFF, 0x00, 0x00) # Red text

                # Example 2: Add a shape (e.g., a simple rectangle outline)
                # left = prs.slide_width / 2
                # top = prs.slide_height / 2
                # width = prs.slide_width / 8
                # height = prs.slide_height / 8
                # shape = slide.shapes.add_shape(1, left, top, width, height) # 1 for rectangle
                # shape.fill.background() # No fill
                # shape.line.color.rgb = (0x00, 0x00, 0xFF) # Blue line
                # shape.line.width = 10000 # 1 point

                # Example 3: Add an image (requires an image file on the server)
                # from pptx.util import Inches
                # img_path = 'path/to/your/logo.png' # Make sure this image exists on your server
                # try:
                #     left = Inches(0.5)
                #     top = Inches(0.05)
                #     slide.shapes.add_picture(img_path, left, top, height=Inches(0.75))
                # except FileNotFoundError:
                #     print(f"Warning: Logo image not found at {img_path}")


                # Example 4: Modify a Master Slide (more complex, but powerful for themes)
                # This would typically be done *before* adding slides, or you'd need to
                # iterate through master slides and modify their elements, or apply
                # a different layout to existing slides.
                # prs.slide_masters[0].shapes.add_textbox(...) # This affects new slides using this master

            # --- End Design Logic ---

            # Save the modified presentation to a BytesIO object
            output_pptx = BytesIO()
            prs.save(output_pptx)
            output_pptx.seek(0) # Rewind to the beginning of the stream

            return send_file(
                output_pptx,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name='designed_presentation.pptx'
            )

        except Exception as e:
            return f"Error processing file: {str(e)}", 500

if __name__ == '__main__':
    # Create a 'templates' directory and put index.html there
    # For development, run with `flask run` or `python server.py`
    app.run(debug=True)
